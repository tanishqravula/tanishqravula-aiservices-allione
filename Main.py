import streamlit as st
import pandas as pd
import google.generativeai as genai
import re
import openpyxl
from PIL import Image
import pdf2image
import pytesseract
import fitz
import os
from pytesseract import Output, TesseractError
from functions import convert_pdf_to_txt_pages, convert_pdf_to_txt_file, save_pages, displayPDF, images_to_txt
from selenium.webdriver.chrome.options import Options
import requests
import PyPDF2 
from docx import Document
from pptx import Presentation
import io
from bs4 import BeautifulSoup
import textwrap
from youtube_transcript_api import YouTubeTranscriptApi
import speech_recognition as sr
from io import StringIO
from io import BytesIO
import html2text
import docx
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.support.wait import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium import webdriver
from bs4 import BeautifulSoup
import asyncio
from dotenv import load_dotenv
#import doc2txt
load_dotenv()

#GOOGLE_API_KEY='AIzaSyAE8nQWFGH1VPC8KRuazzd8SuVQVINq-44'

st.set_page_config(
    page_title="Tanishq AI Chat",
    page_icon="DALL·E 2024-01-15 00.41.44 - A cartoon image of a software developer surrounded by various AI tools. The character is whimsical and colorful, with exaggerated features typical of .png",
    layout="wide",
)
# Path: Main.py

#HEADER
st.markdown('''Powered by Tanishq ''', unsafe_allow_html=True)

image = Image.open("DALL·E 2024-01-15 00.41.44 - A cartoon image of a software developer surrounded by various AI tools. The character is whimsical and colorful, with exaggerated features typical of .png")

st.image(image, caption='ACCESS MY AI PRODUCTS', use_column_width=True)
st.caption("By Tanishq Ravula")

#------------------------------------------------------------
#LANGUAGE
langcols = st.columns([0.2,0.8])
with langcols[0]:
  lang = st.selectbox('Select your language',
  ('Español', 'English', 'Français', 'Deutsch',
  'Italiano', 'Português', 'Polski', 'Nederlands',
  'Русский', '日本語', '한국어', '中文', 'العربية',
  'हिन्दी', 'Türkçe', 'Tiếng Việt', 'Bahasa Indonesia',
  'ภาษาไทย', 'Română', 'Ελληνικά', 'Magyar', 'Čeština',
  'Svenska', 'Norsk', 'Suomi', 'Dansk', 'हिन्दी', 'हिन् '),index=1)

if 'lang' not in st.session_state:
    st.session_state.lang = lang
st.divider()
st.write("For chatting with large or very large PDF Files containing texts use this link:https://tanishqravulachatwithpdfs-ddkmymngr3npso2tpnzuqf.streamlit.app/")
st.write("For chatting with medium sized to somewhat large PDF Files containing images and texts use this link:https://tanishqravulachatwithpdfscontainingsomeimagesandmediumtexts-vc.streamlit.app/")
st.write("For chatting with large or very large CSV or EXCEL Files use this link:https://tanishqravulachatwithlargecsvsandexcels-ymgxpc5ebaapxa8yav4ceh.streamlit.app/")
st.write("For converting your image containing data in the format of table,use this link:https://tanishqravulaimagetoexcelconvertor-ybeofwdwx5vdiajgkhnj3b.streamlit.app/")


#------------------------------------------------------------
#FUNCTIONS
def extract_graphviz_info(text: str) -> list[str]:
  """
  The function extract_graphviz_info takes in a text and returns a list of graphviz code blocks found in the text.

  :param text: The text parameter is a string that contains the text from which you want to extract Graphviz information
  :return: a list of strings that contain either the word "graph" or "digraph". These strings are extracted from the input
  text.
  """

  graphviz_info  = text.split('```')

  return [graph for graph in graphviz_info if ('graph' in graph or 'digraph' in graph) and ('{' in graph and '}' in graph)]

def append_message(message: dict) -> None:
    """
    The function appends a message to a chat session.

    :param message: The message parameter is a dictionary that represents a chat message. It typically contains
    information such as the user who sent the message and the content of the message
    :type message: dict
    :return: The function is not returning anything.
    """
    st.session_state.chat_session.append({'user': message})
    return

@st.cache_data()
def load_model() -> genai.GenerativeModel:
    """
    The function load_model() returns an instance of the genai.GenerativeModel class initialized with the model name
    'gemini-pro'.
    :return: an instance of the genai.GenerativeModel class.
    """
    model = genai.GenerativeModel('gemini-1.5-pro')
    return model

@st.cache_data()
def load_modelvision() -> genai.GenerativeModel:
    """
    The function load_modelvision loads a generative model for vision tasks using the gemini-pro-vision model.
    :return: an instance of the genai.GenerativeModel class.
    """
    model = genai.GenerativeModel('gemini-1.5-flash')
    return model
def to_markdown(text):
    text = text.replace('•', '  *')
    return textwrap.indent(text, '> ', predicate=lambda _: True)

# Function to generate content using GenAI model
def generate_content(model_type, content):
    model = genai.GenerativeModel(model_type)
    response = model.generate_content(content)
    return response.text
def transcribe_video(video_url):
    recognizer = sr.Recognizer()
    audio = None

    with sr.AudioFile(video_url) as source:
        audio = recognizer.record(source)

    return recognizer.recognize_google(audio)
def extract_text_from_docx(docx_file):
    doc = docx.Document(docx_file)
    text = ""

    # Extract text from paragraphs and tables
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text += cell.text + "\t"
            text += "\n"


    return text
def extract_text_from_image(image_bytes):
    image = Image.open(BytesIO(image_bytes))
    text = pytesseract.image_to_string(image)
    return text
def extract_text_from_ppt(ppt_path):
    presentation = Presentation(ppt_path)
    text_content = ""

    for slide_number, slide in enumerate(presentation.slides):
        text_content += f"\n\nSlide {slide_number + 1}:\n"

        for shape_number, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                text_content += shape.text + "\n"

            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                image_text = extract_text_from_image(image_bytes)
                text_content += f"Text from Image {shape_number + 1}:\n{image_text}\n"

    return text_content
def extract_text_from_website(url):
    try:
        response = requests.get(url)
        response.raise_for_status()  # Raise HTTPError for bad requests
        soup = BeautifulSoup(response.text, 'html.parser')
        text_content = soup.get_text()

        # Extract text from tables
        tables = soup.find_all("table")
        table_text = "\n\n".join([pd.read_html(str(table))[0].to_string(index=False) for table in tables])


        return f"{text_content}\n\nTable Text:\n{table_text}\n"
    except Exception as e:
        return "XYZ#&^^@^%@hx"
def extract_text_from_images_on_website(images):
    extracted_text = ""
    for image in images:
        try:
            # Convert the image to an array
            image_np = np.array(image)

            # Convert RGB to BGR (OpenCV uses BGR format)
            image_cv2 = cv2.cvtColor(image_np, cv2.COLOR_RGB2BGR)

            # Use Tesseract to extract text
            text = pytesseract.image_to_string(image_cv2, lang='eng')
            extracted_text += f"Text from Image: {text}\n"

        except Exception as e:
            extracted_text += f"Error extracting text from image: {str(e)}\n"

    return extracted_text
def extract_content_with_selenium(url):
    try:
        # Configure Chrome options for running in headless mode
        chrome_options = Options()
        chrome_options.add_argument("--headless")
        chrome_options.add_argument('--disable-gpu')
        chrome_options.add_argument('--no-sandbox')
        chrome_options.add_argument("--headless")
        chrome_options.add_argument("--no-sandbox")
        chrome_options.add_argument("--disable-dev-shm-usage")
        chrome_options.add_argument("--disable-gpu")
        chrome_options.add_argument("--disable-features=NetworkService")
        chrome_options.add_argument("--window-size=1920x1080")
        chrome_options.add_argument("--disable-features=VizDisplayCompositor")
        chrome_options.add_argument('--ignore-certificate-errors')
 

        # Create the driver with the options
        driver = webdriver.Chrome(options=chrome_options)

        # Load the page with Selenium
        driver.get(url)

        # Wait up to 10 seconds for the page to load
        # Wait for the page to finish loading all JavaScript
        wait = WebDriverWait(driver, 10)
        wait.until(EC.presence_of_element_located((By.XPATH, "//body[not(@class='loading')]")))

        # Get the HTML of the page
        html = driver.page_source

        # Close the WebDriver
        driver.quit()

        # Parse HTML content using BeautifulSoup
        soup = BeautifulSoup(html, 'html.parser')

        # Extract desired content here
        # For example, let's extract all text from paragraphs
        paragraphs = soup.find_all('p')
        text_content = '\n'.join([p.get_text() for p in paragraphs])
        para_content = soup.get_text(separator='\n')
        #text_content1=soup.get_text(separator='\n')


        # Extract table content
        tables = soup.find_all('table')
        table_content = ""
        for table in tables:
            rows = table.find_all('tr')
            for row in rows:
                cols = row.find_all(['th', 'td'])
                row_text = '\t'.join([col.get_text() for col in cols])
                table_content += row_text + '\n'

        # Close the browser
        #driver.quit()

        return text_content, table_content, para_content
    except Exception as e:
        st.error(f"Error extracting content from the website with Selenium: {e}")
        return "","",""
def generate_gemini(model_type, content):
    try:
        #content=str(content)
        model = load_model()  # Ensure model is loaded
        response = model.generate_content(content)
        if hasattr(response, 'text'):
            return response.text
        else:
            # Convert the response object to string and return
            return str(response)
    except ValueError as e:
        error_message = f"An error occurred while generating content: {str(e)}"
        # Log the error message or handle it appropriately
        print(error_message)
        return None





#------------------------------------------------------------
#CONFIGURATION
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

model = load_model()

vision = load_modelvision()

if 'chat' not in st.session_state:
    st.session_state.chat = model.start_chat(history=[])

if 'chat_session' not in st.session_state:
    st.session_state.chat_session = []

#st.session_state.chat_session

#------------------------------------------------------------
#CHAT

if 'messages' not in st.session_state:
    st.session_state.messages = []

if 'welcome' not in st.session_state or lang != st.session_state.lang:
    st.session_state.lang = lang
    welcome  = 'Welcome! I am a chatbot in a chat application created using Streamlit and Python. I can help you with a variety of tasks such as describing images, answering questions, reading text files, reading tables, generating graphs with Graphviz,reading documents,chatting with websites,chatting with youtube urls and more. How can I assist you today?'

    st.session_state.welcome =welcome

    with st.chat_message('ai'):
        st.write(st.session_state.welcome)
else:
    with st.chat_message('ai'):
        st.write(st.session_state.welcome)

if len(st.session_state.chat_session) > 0:
    count = 0
    for message in st.session_state.chat_session:

        if message['user']['role'] == 'model':
            with st.chat_message('ai'):
                st.write(message['user']['parts'])
                graphs = extract_graphviz_info(message['user']['parts'])
                if len(graphs) > 0:
                    for graph in graphs:
                        st.graphviz_chart(graph,use_container_width=False)
                        if lang == 'Español':
                          view = "Ver texto"
                        else:
                          view = "View text"
                        with st.expander(view):
                          st.code(graph, language='dot')
        else:
            with st.chat_message('user'):
                st.write(message['user']['parts'][0])
                if len(message['user']['parts']) > 1:
                    st.image(message['user']['parts'][1], width=200)
        count += 1



#st.session_state.chat.history

cols=st.columns(7)

with cols[0]:
    if lang == 'Español':
      image_atachment = st.toggle("Adjuntar imagen", value=False, help="Activa este modo para adjuntar una imagen y que el chatbot pueda leerla")
    else:
      image_atachment = st.toggle("Attach image", value=False, help="Activate this mode to attach an image and let the chatbot read it")

with cols[1]:
    if lang == 'Español':
      txt_atachment = st.toggle("Adjuntar archivo de texto", value=False, help="Activa este modo para adjuntar un archivo de texto y que el chatbot pueda leerlo")
    else:
      txt_atachment = st.toggle("Attach text file", value=False, help="Activate this mode to attach a text file and let the chatbot read it")
with cols[2]:
    if lang == 'Español':
      csv_excel_atachment = st.toggle("Adjuntar CSV o Excel", value=False, help="Activa este modo para adjuntar un archivo CSV o Excel y que el chatbot pueda leerlo")
    else:
      csv_excel_atachment = st.toggle("Attach CSV or Excel", value=False, help="Activate this mode to attach a CSV or Excel file and let the chatbot read it")
with cols[3]:
    if lang == 'Español':
      graphviz_mode = st.toggle("Modo graphviz", value=False, help="Activa este modo para generar un grafo con graphviz en .dot a partir de tu mensaje")
    else:
      graphviz_mode = st.toggle("Graphviz mode", value=False, help="Activate this mode to generate a graph with graphviz in .dot from your message")
with cols[4]:
    if lang == 'Español':
        doc_atachment = st.toggle("Adjuntar PDF, PPTX, DOCX", value=False, help="Activa este modo para adjuntar un archivo PDF, PPT o DOCX y que el chatbot pueda leerlo")
    else:

        doc_atachment = st.toggle("Attach PDF, PPTX, DOCX", value=False, help="Activate this mode to attach a PDF, PPT, or DOCX file and let the chatbot read it")
with cols[5]:
    if lang == 'Español':
        website_chat = st.toggle("Chatear con sitios web", value=False,
                                 help="Activa este modo para chatear con un sitio web y resumir su contenido")
    else:
        website_chat = st.toggle("Chat with websites", value=False,
                                 help="Activate this mode to chat with a website and summarize its content")
with cols[6]:
    if lang == 'Español':
        youtube_chat = st.toggle("Chatear con youtube urls", value=False,
                                 help="Activa este modo para chatear con un sitio yotube resumir su contenido")
    else:
        youtube_chat = st.toggle("Chat with youtube urls", value=False,
                                 help="Activate this mode to chat with a youtube url or youtube video and chat with it") 


if image_atachment:
    if lang == 'Español':
      image = st.file_uploader("Sube tu imagen", type=['png', 'jpg', 'jpeg'])
      url = st.text_input("O pega la url de tu imagen")
    else:
      image = st.file_uploader("Upload your image", type=['png', 'jpg', 'jpeg'])
      url = st.text_input("Or paste your image url")
else:
    image = None
    url = ''



if txt_atachment:
    if lang == 'Español':
      txtattachment = st.file_uploader("Sube tu archivo de texto", type=['txt'])
    else:
      txtattachment = st.file_uploader("Upload your text file", type=['txt'])
else:
    txtattachment = None
if doc_atachment:
    if lang == 'Español':
        docattachment = st.file_uploader("Sube tu archivo PDF, PPT, DOCX", type=['pdf', 'pptx', 'docx'])
    else:
        docattachment = st.file_uploader("Upload your PDF, PPT, DOCX file", type=['pdf', 'pptx', 'docx'])
else:
    docattachment = None

if csv_excel_atachment:
    if lang == 'Español':
      csvexcelattachment = st.file_uploader("Sube tu archivo CSV o Excel", type=['csv', 'xlsx'])
    else:
      csvexcelattachment = st.file_uploader("Upload your CSV or Excel file", type=['csv', 'xlsx'])
else:
    csvexcelattachment = None
if youtube_chat:
    if lang == 'Español':
        youtube_url = st.text_input("Introduce la URL del sitio youtube:")
    else:
        youtube_url = st.text_input("Enter the URL of the youtube:")
    if youtube_url:
        video_id = youtube_url.split("=")[1]  # Extracting the video ID

        try:
            video_transcript = YouTubeTranscriptApi.get_transcript(video_id)
            if video_transcript:
                full_transcript = ' '.join([x['text'] for x in video_transcript])
            else:
                # If no transcript available, attempt to transcribe the speech
                # Assuming the video has no transcript available
                st.warning("No transcript available. Attempting to transcribe the speech...")
                full_transcript = transcribe_video(youtube_url)

            #st.subheader("Video Transcript")
            #st.write(full_transcript)

            #st.subheader("Summary")
            #summarized_text = summarize_text(full_transcript)
            #st.write(summarized_text)

        except Exception as e:
            st.error(f"An error occurred: {e}")


if website_chat:
    if lang == 'Español':
        website_url = st.text_input("Introduce la URL del sitio web:")
    else:
        website_url = st.text_input("Enter the URL of the website:")

    if website_url:
        try:
            website_text=''
            text_content, table_content,para_content = extract_content_with_selenium(website_url)
            website_text+=text_content
            website_text+=table_content
            website_text+=para_content
            #if(extract_text_from_website(website_url)!='XYZ#&^^@^%@hx'):
                #website_text+=extract_text_from_website(website_url)
            content=f'summarise this content briefly:{website_text} without missing even one word from the text fetched from information:{website_text} and complete the whole generated content'
            content1=f'organize the content: {website_text} into  tables '
            #result = generate_gemini("gemini-pro", content)
            #result1=generate_gemini("gemini-pro",content1)
            #result = generate_gemini("gemini-pro", content)
            #result1=generate_gemini("gemini-pro",content1)
                
                


            # Summarize the text if needed
            # (You can use a summarization library or method here)

            # Display the summarized text in the chat
            with st.chat_message('user'):
                st.write(f"Content: {website_url}")
            with st.chat_message('model'):
                #if(generate_gemini("gemini-pro",website_text)=='' and generate_gemini("gemini-pro",content1)==''):
                st.write(f'Extracted content from website:{website_text}')  
                st.markdown(to_markdown(generate_gemini("gemini-pro", content)))
                st.markdown(to_markdown(generate_gemini("gemini-pro",content1)))
                    #st.write(f'Extracted content from website:{website_text}')
                



        except Exception as e:
            with st.chat_message('model'):
                st.write(f"The website  does not allow to collect and fetch information according to the websites privacy and confidential information.You can try another website urls {str(e)}")
if lang == 'Español':
  prompt = st.chat_input("Escribe tu mensaje")
else:
  prompt = st.chat_input("Write your message")

if prompt:
    txt = ''
    if txtattachment:
        txt = txtattachment.getvalue().decode("utf-8")
        if lang == 'Español':
          txt = '   Archivo de texto: \n' + txt
        else:
          txt = '   Text file: \n' + txt
    if docattachment:
        path =docattachment.read()
        file_extension = docattachment.name.split('.')[-1].lower()
        try:
            if file_extension == 'pdf':
                
                if docattachment is not None:
                    doc = fitz.open(stream=path, filetype="pdf")
                    doc_content = ""
                    for page_num in range(len(doc)):
                        page = doc.load_page(page_num)
                        doc_content += page.get_text()
                    texts, nbPages = images_to_txt(path, 'eng')
                    text_data_f = "\n\n".join(texts)
                    #text_data_text, nbPages_text = convert_pdf_to_txt_file(docattachment)
                    doc_content+=text_data_f
            elif file_extension == 'docx'or file_extension=='doc':
                docx_content = extract_text_from_docx(docattachment)

            elif file_extension == 'ppt' or file_extension == 'pptx':
                ppt_content =  extract_text_from_ppt(docattachment)
            else:
                raise ValueError(f'Unsupported file format: {file_extension}')

        except Exception as e:
            st.error(f"Error processing {file_extension} file: {e}")
        # Add more specific handling/logging if needed

        else:
            if file_extension == 'pdf':
                if lang == 'Español':
                    txt += f'   Archivo adjunto (PDF): \n{doc_content}'
                else:
                    txt += f'   Attached file (PDF): \n{doc_content}'
            elif file_extension == 'docx':
                if lang == 'Español':
                    txt += f'   Archivo adjunto (DOCX): \n{docx_content}'
                else:
                    txt += f'   Attached file (DOCX): \n{docx_content}'
            elif file_extension == 'ppt' or file_extension == 'pptx':
                if lang == 'Español':
                    txt += f'   Archivo adjunto (PPT): \n{ppt_content}'
                else:
                    txt += f'   Attached file (PPT): \n{ppt_content}'
    




    if csvexcelattachment:
        file_name = csvexcelattachment.name
        file_extension = file_name.split('.')[-1].lower()

        if file_extension == 'csv':
             df = pd.read_csv(csvexcelattachment)
             txt += '   Dataframe: \n' + df.to_string()
        elif file_extension in ['xlsx', 'xls']:
            wb = openpyxl.load_workbook(csvexcelattachment)
            sheet = wb.active
            for row in sheet.iter_rows():
            # extract each cell value
                for cell in row:
                    txt += str(cell.value)
        else:
            st.error("Unsupported file type. Please upload a CSV or Excel file.")
    if website_chat:
        txt=website_text
        prmt  = {'role': 'user', 'parts':[prompt+txt]}
    else:
        st.write('')
        prmt  = {'role': 'user', 'parts':[prompt+txt]}
    


    if graphviz_mode:
        if lang == 'Español':
          txt += '   Genera un grafo con graphviz en .dot \n'
        else:
          txt += '   Generate a graph with graphviz in .dot \n'
    if youtube_chat:
        txt=full_transcript 
        prmt  = {'role': 'user', 'parts':[prompt+txt]}
    else:
        st.write('')
        prmt  = {'role': 'user', 'parts':[prompt+txt]}

    if len(txt) > 9000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000009999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999:
        txt = txt[:9000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000000009999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999999] + '...'
    if image or url != '':
        if url != '':
            img = Image.open(requests.get(url, stream=True).raw)
        else:
            img = Image.open(image)
        prmt  = {'role': 'user', 'parts':[prompt+txt, img]}
    else:
        prmt  = {'role': 'user', 'parts':[prompt+txt]}

    append_message(prmt)

    if lang == 'Español':
      spinertxt = 'Espera un momento, estoy pensando...'
    else:
      spinertxt = 'Wait a moment, I am thinking...'
    with st.spinner(spinertxt):
        if len(prmt['parts']) > 0:
            if image_atachment:
                response = vision.generate_content(prmt['parts'],stream=True)
            else:
                response = vision.generate_content(prmt['parts'],stream=True)
            #response.resolve()
        else:
            response = st.session_state.chat.send_message(prmt['parts'][0])

        try:
          append_message({'role': 'model', 'parts':response.text})
        except Exception as e:
          append_message({'role': 'model', 'parts':f'{type(e).name}: {e}'})


        st.rerun()



#st.session_state.chat_session
